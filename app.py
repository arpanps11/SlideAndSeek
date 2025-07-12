from flask import Flask, render_template, request, redirect, url_for, session, flash, send_file
import sqlite3
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from datetime import datetime
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)
app.secret_key = 'supersecretkey_root'

DB_FILE = 'songs.db'
PASSWORD = 'ebccni@2025'

def get_db_connection():
    conn = sqlite3.connect(DB_FILE)
    conn.row_factory = sqlite3.Row
    return conn

# ---------- Routes ----------

@app.route('/')
def home():
    return render_template('home.html')

@app.route('/verify', methods=['GET', 'POST'])
def verify():
    error = None
    next_page = request.args.get('next', '/')
    if request.method == 'POST':
        password = request.form.get('password')
        if password == PASSWORD:
            session['authenticated'] = True
            return redirect(next_page)
        else:
            error = "Incorrect password. Try again."
    return render_template('verify.html', error=error, next=next_page)

# ---------- Song Management Routes (add/edit/search) ----------
@app.route('/add', methods=['GET', 'POST'])
def add_song():
    if not session.get('authenticated'):
        return redirect(url_for('verify', next='/add'))

    error = None

    if request.method == 'POST':
        title = request.form.get('title', '').strip()
        lyrics = request.form.get('lyrics', '').strip()
        song_number = request.form.get('song_number')
        page_number = request.form.get('page_number')
        key_root = request.form.get('key_root')
        key_type = request.form.get('key_type')

        if not title or not lyrics:
            error = "Song Title and Lyrics are required."
        else:
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("SELECT * FROM songs WHERE LOWER(title) = LOWER(?)", (title.lower(),))
            existing = cursor.fetchone()

            if existing:
                error = "Song title already exists."
            else:
                cursor.execute("""
                    INSERT INTO songs (title, lyrics, song_number, page_number, key_root, key_type)
                    VALUES (?, ?, ?, ?, ?, ?)
                """, (title, lyrics, song_number, page_number, key_root, key_type))
                conn.commit()
                flash("Song added successfully.")
            conn.close()

    return render_template('add.html', error=error)

@app.route('/edit/<int:song_id>', methods=['GET', 'POST'])
def edit_song(song_id):
    if not session.get('authenticated'):
        return redirect(url_for('verify', next=f'/edit/{song_id}'))

    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT * FROM songs WHERE id = ?", (song_id,))
    song = cursor.fetchone()

    if not song:
        conn.close()
        return "Song not found", 404

    error = None

    if request.method == 'POST':
        title = request.form.get('title', '').strip()
        lyrics = request.form.get('lyrics', '').strip()
        song_number = request.form.get('song_number')
        page_number = request.form.get('page_number')
        key_root = request.form.get('key_root')
        key_type = request.form.get('key_type')

        if not title or not lyrics:
            error = "Title and lyrics are required."
        else:
            cursor.execute("SELECT id FROM songs WHERE LOWER(title) = LOWER(?) AND id != ?", (title.lower(), song_id))
            duplicate = cursor.fetchone()
            if duplicate:
                error = "Another song with this title already exists."
            else:
                cursor.execute("""
                    UPDATE songs
                    SET title = ?, lyrics = ?, song_number = ?, page_number = ?, key_root = ?, key_type = ?
                    WHERE id = ?
                """, (title, lyrics, song_number, page_number, key_root, key_type, song_id))
                conn.commit()
                conn.close()
                flash("Song updated successfully.")
                return redirect(url_for('search', query=title))
    else:
        conn.close()

    return render_template('edit.html', song=song, error=error)

@app.route('/search', methods=['GET', 'POST'])
def search():
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT title FROM songs")
    all_songs = [{'title': row[0]} for row in cursor.fetchall()]

    query = request.args.get('query', '').strip().lower()
    results = []
    searched = bool(query)

    if request.method == 'POST':
        query = request.form.get('query', '').strip().lower()
        searched = True

    if searched:
        cursor.execute("""
            SELECT title, lyrics, key_root, key_type, song_number, id FROM songs
            WHERE LOWER(title) LIKE ? 
               OR LOWER(lyrics) LIKE ? 
               OR LOWER(COALESCE(key_root, '') || ' ' || COALESCE(key_type, '')) LIKE ?
        """, (f'%{query}%', f'%{query}%', f'%{query}%'))
        results = cursor.fetchall()

    conn.close()
    return render_template('search.html', results=results, searched=searched, query=query, all_songs=all_songs)

# ---------- Responsive Reading ----------
@app.route('/rr_home')
def rr_home():
    return render_template('rr_home.html')

@app.route('/rr_add', methods=['GET', 'POST'])
def rr_add():
    if not session.get('authenticated'):
        return redirect(url_for('verify', next='/rr_add'))

    error = None

    if request.method == 'POST':
        rr_number = request.form.get('rr_number')
        psalm_number = request.form.get('psalm_number')
        page_number = request.form.get('page_number')
        title = request.form.get('title', '').strip()
        content = request.form.get('content', '').strip()

        if not all([rr_number, psalm_number, page_number, title, content]):
            error = "All fields are required."
        else:
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("""
                INSERT INTO responsive_readings (rr_number, psalm_number, page_number, title, content)
                VALUES (?, ?, ?, ?, ?)
            """, (rr_number, psalm_number, page_number, title, content))
            conn.commit()
            conn.close()
            flash("Responsive Reading added successfully.")
            return redirect(url_for('rr_add'))

    return render_template('rr_add.html', error=error)

@app.route('/rr_search', methods=['GET', 'POST'])
def rr_search():
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT id, rr_number, psalm_number FROM responsive_readings ORDER BY rr_number ASC")
    rr_list = cursor.fetchall()

    rr = None
    selected_id = None
    searched = False

    if request.method == 'POST' or request.args.get('selected_id'):
        selected_id = request.form.get('rr_id') or request.args.get('selected_id')
        searched = True
        if selected_id:
            cursor.execute("SELECT * FROM responsive_readings WHERE id = ?", (selected_id,))
            rr = cursor.fetchone()

    conn.close()
    return render_template('rr_search.html', rr_list=rr_list, rr=rr, selected_id=selected_id, searched=searched)

@app.route('/rr_edit/<int:rr_id>', methods=['GET', 'POST'])
def rr_edit(rr_id):
    if not session.get('authenticated'):
        return redirect(url_for('verify', next=f'/rr_edit/{rr_id}'))

    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT * FROM responsive_readings WHERE id = ?", (rr_id,))
    rr = cursor.fetchone()

    if not rr:
        conn.close()
        return "Responsive Reading not found", 404

    error = None

    if request.method == 'POST':
        rr_number = request.form.get('rr_number')
        psalm_number = request.form.get('psalm_number')
        page_number = request.form.get('page_number')
        title = request.form.get('title', '').strip()
        content = request.form.get('content', '').strip()

        if not all([rr_number, psalm_number, page_number, title, content]):
            error = "All fields are required."
        else:
            cursor.execute("""
                UPDATE responsive_readings
                SET rr_number = ?, psalm_number = ?, page_number = ?, title = ?, content = ?
                WHERE id = ?
            """, (rr_number, psalm_number, page_number, title, content, rr_id))
            conn.commit()
            conn.close()
            flash("Responsive Reading updated successfully.")
            return redirect(url_for('rr_search', selected_id=rr_id))

    conn.close()
    return render_template('rr_edit.html', rr=rr, error=error)

# ---------- Slide Generator ----------
@app.route('/generate', methods=['GET', 'POST'])
def generate():
    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute("SELECT * FROM songs ORDER BY title ASC")
    songs = cursor.fetchall()
    cursor.execute("SELECT * FROM responsive_readings ORDER BY rr_number ASC")
    rrs = cursor.fetchall()
    conn.close()

    if request.method == 'POST':
        def get_song_by_id(song_id):
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("SELECT * FROM songs WHERE id = ?", (song_id,))
            song = cursor.fetchone()
            conn.close()
            return song

        def get_rr_by_id(rr_id):
            conn = get_db_connection()
            cursor = conn.cursor()
            cursor.execute("SELECT * FROM responsive_readings WHERE id = ?", (rr_id,))
            rr = cursor.fetchone()
            conn.close()
            return rr

        section_order = request.form.getlist('section_order[]')
        prs = Presentation()

        for section in section_order:
            if section == 'welcome_section':
                welcome_text = request.form.get('welcome_text', '').strip()
                if welcome_text:
                    add_content_slide(prs, welcome_text)

            elif section == 'praise_section':
                praise_ids = request.form.getlist('praise_ids[]')
                if praise_ids:
                    add_title_slide(prs, "Praise and Worship")
                for song_id in praise_ids:
                    song = get_song_by_id(int(song_id))
                    if song:
                        add_title_slide(prs, song['title'])
                        add_song_slides(prs, song)

            elif section == 'opening_section':
                song_id = request.form.get('opening_id')
                if song_id:

                    song = get_song_by_id(int(song_id))
                    if song:
                        meta = compose_song_meta(song)
                        add_title_slide(prs, 'Opening Hymn', f"{song['title']} ({meta})" if meta else song['title'])
                        add_song_slides(prs, song)

            elif section == 'offertory_section':
                song_id = request.form.get('offertory_id')
                if song_id:
                    song = get_song_by_id(int(song_id))
                    if song:
                        meta = compose_song_meta(song)
                        add_title_slide(prs, 'Offertory Hymn', f"{song['title']} ({meta})" if meta else song['title'])
                        add_song_slides(prs, song)

            elif section == 'communion_section':
                song_id = request.form.get('communion_id')
                if song_id:
                    song = get_song_by_id(int(song_id))
                    if song:
                        meta = compose_song_meta(song)
                        add_title_slide(prs, 'Communion Hymn', f"{song['title']} ({meta})" if meta else song['title'])
                        add_song_slides(prs, song)

            elif section == 'rr_section':
                rr_id = request.form.get('rr_id')
                if rr_id:
                    rr = get_rr_by_id(int(rr_id))
                    if rr:
                        rr_intro = f"Responsive Reading {rr['rr_number']}\nPsalm {rr['psalm_number']}\n{rr['title']}"
                        rr_page = rr['page_number']
                        if rr_page:
                            rr_intro += f"\nPage {rr_page}"
                        add_content_slide(prs, rr_intro, font_size=40)

                        content = rr['content'].replace('\r\n', '\n').replace('_x000D_', '').strip()
                        verses = [v.strip() for v in content.split('\n\n') if v.strip()]
                        for verse in verses:
                            add_content_slide(prs, verse, font_size=32)

            elif section == 'extras_section':
                extras_title = request.form.get('extras_title', 'Extra Songs').strip()
                extras_ids = request.form.getlist('extras_ids[]')
                for song_id in extras_ids:
                    song = get_song_by_id(int(song_id))
                    if song:
                        add_title_slide(prs, extras_title, song['title'])
                        add_song_slides(prs, song)

            elif section == 'doxology_section':
                add_title_slide(prs, 'Doxology')
                benedict_lines = [
                    "Praise God From Who All Blessings Flow,",
                    "Praise Him All Creatures Here Below,",
                    "Praise Him Above Ye Heavenly Hosts",
                    "Praise Father, Son and Holy Ghost",
                    "Amen"
                ]
                add_content_slide(prs, '\n'.join(benedict_lines))

        ppt_io = BytesIO()
        prs.save(ppt_io)
        ppt_io.seek(0)
        filename = f"Worship_{datetime.now().strftime('%d_%m_%Y')}.pptx"
        return send_file(
            ppt_io,
            as_attachment=True,
            download_name=filename,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )

    # ✅ Convert songs and RRs to JSON-serializable dictionaries
    songs_dicts = [dict(song) for song in songs]
    rrs_dicts = [dict(rr) for rr in rrs]

    return render_template('generate.html', songs=songs_dicts, rrs=rrs_dicts)

# ---------- Utility Functions ----------
def add_title_slide(ppt, title, subtitle=None):
    slide_layout = ppt.slide_layouts[0]
    slide = ppt.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    if subtitle:
        slide.placeholders[1].text = subtitle

def add_content_slide(prs, text, font_size=32):
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(9)
    height = Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(font_size)
    p.font.bold = True

def add_song_slides(ppt, song, include_title_slide=False):
    if include_title_slide:
        meta = compose_song_meta(song)
        add_title_slide(ppt, song['title'], meta)
    stanzas = split_into_paragraphs(song['lyrics'])
    for stanza in stanzas:
        stanza_clean = stanza.replace('_x000D_', '').strip()
        if stanza_clean:
            add_content_slide(ppt, stanza_clean)

def split_into_paragraphs(text):
    return [p.strip() for p in text.replace('\r', '').split('\n\n') if p.strip()]

def compose_song_meta(song):
    parts = []
    if song['song_number']:
        parts.append(f"Song {song['song_number']}")
    if song['page_number']:
        parts.append(f"Page {song['page_number']}")
    return " | ".join(parts) if parts else ""

# ---------- Initialize DB ----------
if __name__ == '__main__':
    if not os.path.exists(DB_FILE):
        conn = get_db_connection()
        cursor = conn.cursor()
        cursor.execute('''
            CREATE TABLE songs (
                id INTEGER PRIMARY KEY AUTOINCREMENT,
                title TEXT NOT NULL,
                lyrics TEXT NOT NULL,
                song_number TEXT,
                page_number TEXT,
                key_root TEXT,
                key_type TEXT
            )
        ''')
        conn.commit()
        conn.close()

    conn = get_db_connection()
    cursor = conn.cursor()
    cursor.execute('''
        CREATE TABLE IF NOT EXISTS responsive_readings (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            rr_number INTEGER NOT NULL,
            psalm_number INTEGER NOT NULL,
            page_number TEXT NOT NULL,
            title TEXT NOT NULL,
            content TEXT NOT NULL
        )
    ''')
    conn.commit()
    conn.close()

    app.run(debug=True)
